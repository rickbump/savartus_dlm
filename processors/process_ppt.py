
import os
from pptx import Presentation


def process_ppt(file_path):
    metadata = {}
    try:
        # Basic file metadata
        file_metadata = os.stat(file_path)
        metadata['file_size'] = file_metadata.st_size
        metadata['file_modified'] = file_metadata.st_mtime

        # Load PowerPoint presentation
        presentation = Presentation(file_path)

        # Presentation properties (core properties)
        metadata['slide_count'] = len(presentation.slides)

        # Process metadata for each slide
        slides_metadata = []
        for slide_index, slide in enumerate(presentation.slides):
            slide_info = {
                'slide_index': slide_index + 1,
                'shape_count': len(slide.shapes),
                'table_count': sum(1 for shape in slide.shapes if shape.has_table)
            }

            # Process tables within each slide
            tables_metadata = []
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    table_info = {
                        'row_count': len(table.rows),
                        'column_count': len(table.columns),
                    }
                    tables_metadata.append(table_info)
            slide_info['tables'] = tables_metadata

            # Store slide metadata
            slides_metadata.append(slide_info)
        metadata['slides'] = slides_metadata

    except Exception as e:
        metadata['ppt_error'] = str(e)

    return metadata